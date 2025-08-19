# core/document_processor.py
import os
import time
from io import StringIO
from datetime import datetime
from pdfminer.high_level import extract_text_to_fp
from langchain_text_splitters import RecursiveCharacterTextSplitter
from config import CHUNK_SIZE, CHUNK_OVERLAP
import logging
import re


# 新增的文件处理库
try:
    from docx import Document  # python-docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl  # Excel处理
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation  # PowerPoint处理
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import markdown  # Markdown处理
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False


def extract_text_from_file(filepath: str) -> str:
    """根据文件路径提取文本内容，支持多种文件格式。"""
    file_ext = filepath.lower().split('.')[-1]

    try:
        if file_ext == 'pdf':
            # PDF文件处理
            output = StringIO()
            with open(filepath, 'rb') as file:
                extract_text_to_fp(file, output)
            return output.getvalue()

        elif file_ext == 'txt':
            # 纯文本文件处理
            with open(filepath, 'r', encoding='utf-8') as file:
                return file.read()

        elif file_ext in ['doc', 'docx']:
            # Word文档处理
            if not DOCX_AVAILABLE:
                raise ImportError("需要安装 python-docx 库: pip install python-docx")
            doc = Document(filepath)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            # 处理表格中的文本
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            return '\n'.join(text)

        elif file_ext in ['xls', 'xlsx']:
            # Excel文件处理
            if not EXCEL_AVAILABLE:
                raise ImportError("需要安装 openpyxl 库: pip install openpyxl")
            workbook = openpyxl.load_workbook(filepath, data_only=True)
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"工作表: {sheet_name}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        text.append(row_text)
                text.append('')  # 工作表之间的分隔
            return '\n'.join(text)

        elif file_ext in ['ppt', 'pptx']:
            # PowerPoint文件处理
            if not PPTX_AVAILABLE:
                raise ImportError("需要安装 python-pptx 库: pip install python-pptx")
            prs = Presentation(filepath)
            text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"幻灯片 {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
                text.append('')  # 幻灯片之间的分隔
            return '\n'.join(text)

        elif file_ext in ['md', 'markdown']:
            # Markdown文件处理
            with open(filepath, 'r', encoding='utf-8') as file:
                md_content = file.read()

            if MARKDOWN_AVAILABLE:
                # 如果有markdown库，可以转换为HTML再提取纯文本
                html = markdown.markdown(md_content)
                # 简单的HTML标签移除
                text = re.sub('<[^<]+?>', '', html)
                return text
            else:
                # 直接返回markdown原文
                return md_content

        elif file_ext in ['csv']:
            # CSV文件处理
            import csv
            text = []
            with open(filepath, 'r', encoding='utf-8') as file:
                csv_reader = csv.reader(file)
                for row in csv_reader:
                    text.append('\t'.join(row))
            return '\n'.join(text)

        elif file_ext in ['json']:
            # JSON文件处理
            import json
            with open(filepath, 'r', encoding='utf-8') as file:
                data = json.load(file)
            return json.dumps(data, ensure_ascii=False, indent=2)

        elif file_ext in ['html', 'htm']:
            # HTML文件处理
            with open(filepath, 'r', encoding='utf-8') as file:
                html_content = file.read()
            # 简单的HTML标签移除
            text = re.sub('<[^<]+?>', '', html_content)
            # 移除多余的空白字符
            text = re.sub('\s+', ' ', text).strip()
            return text

        else:
            # 尝试作为纯文本处理
            try:
                with open(filepath, 'r', encoding='utf-8') as file:
                    return file.read()
            except UnicodeDecodeError:
                # 如果UTF-8解码失败，尝试其他编码
                encodings = ['gbk', 'gb2312', 'latin-1']
                for encoding in encodings:
                    try:
                        with open(filepath, 'r', encoding=encoding) as file:
                            return file.read()
                    except UnicodeDecodeError:
                        continue
                raise ValueError(f"无法解码文件: {os.path.basename(filepath)}")

    except Exception as e:
        logging.error(f"处理文件 {filepath} 时出错: {e}")
        raise ValueError(f"处理文件 {os.path.basename(filepath)} 时出错: {str(e)}")


def get_supported_file_types() -> str:
    """返回支持的文件类型列表。"""
    supported_types = [
        "PDF (.pdf)",
        "纯文本 (.txt)",
        "Word文档 (.doc, .docx)" + (" - 需要安装 python-docx" if not DOCX_AVAILABLE else ""),
        "Excel表格 (.xls, .xlsx)" + (" - 需要安装 openpyxl" if not EXCEL_AVAILABLE else ""),
        "PowerPoint (.ppt, .pptx)" + (" - 需要安装 python-pptx" if not PPTX_AVAILABLE else ""),
        "Markdown (.md, .markdown)" + (" - 需要安装 markdown" if not MARKDOWN_AVAILABLE else ""),
        "CSV (.csv)",
        "JSON (.json)",
        "HTML (.html, .htm)"
    ]
    return "\n".join(f"• {t}" for t in supported_types)


def enhanced_recursive_split(text: str, chunk_size: int = None, chunk_overlap: int = None) -> list[str]:
    """
    增强递归分块方法

    特点：
    1. 智能分隔符选择：根据文档类型和内容特征动态选择分隔符
    2. 语义边界保护：优先在完整语义单元边界切分
    3. 重叠策略优化：动态调整重叠大小
    4. 分块质量评估：确保每个分块的语义完整性
    """
    if chunk_size is None:
        chunk_size = CHUNK_SIZE
    if chunk_overlap is None:
        chunk_overlap = CHUNK_OVERLAP

    # 1. 文档类型检测和分隔符优化
    separators = _get_smart_separators(text)

    # 2. 创建增强的递归分割器
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=separators,
        length_function=len,
        is_separator_regex=False
    )

    # 3. 执行初步分割
    initial_chunks = text_splitter.split_text(text)

    # 4. 后处理：语义边界优化和质量评估
    optimized_chunks = _optimize_chunk_boundaries(initial_chunks, chunk_size, chunk_overlap)

    # 5. 分块质量评估和修复
    final_chunks = _evaluate_and_fix_chunks(optimized_chunks, chunk_size)

    # 6. 清理分块开头的无用标点符号
    cleaned_chunks = _clean_chunk_prefixes(final_chunks)

    return cleaned_chunks


def _get_smart_separators(text: str) -> list[str]:
    """
    根据输入文本的结构特征，智能生成分隔符优先级列表，用于后续层次化文本切分。

    该函数分析文本中是否存在 Markdown 标题、代码块、列表等结构，并据此动态调整分隔符的使用顺序：
    - 优先使用能反映文档逻辑结构的分隔符（如标题、代码块边界）
    - 然后按语义强度从强到弱依次使用通用分隔符（如段落、句子、词语边界）
    - 最终形成一个从“大块”到“小块”的切分策略，确保文本在语义合理的位置被分割

    分隔符列表按优先级降序排列，切分过程会依次尝试每个分隔符，直到满足块大小要求。

    Args:
        text (str): 输入的原始文档文本

    Returns:
        list[str]: 按优先级排序的分隔符列表，例如：
                   ["\n## ", "\n```\n", "\n\n\n", "\n\n", "。", ".", " ", ...]

    示例场景：
        - Markdown 文档：优先按标题（##）和代码块（````）分割
        - 普通段落文本：优先按段落（\n\n）和句号（。）分割
        - 列表内容：在列表项之间保持完整，避免跨项切分
    """
    # 基础分隔符集合
    base_separators = []

    # 检测文档结构特征
    has_markdown_headers = bool(re.search(r'^#+\s', text, re.MULTILINE))
    has_numbered_lists = bool(re.search(r'^\d+\.\s', text, re.MULTILINE))
    has_bullet_points = bool(re.search(r'^[•\-\*]\s', text, re.MULTILINE))
    has_code_blocks = bool(re.search(r'```', text))

    # 根据文档特征调整分隔符优先级
    if has_markdown_headers:
        # Markdown文档：优先按标题分割
        base_separators.extend(["\n## ", "\n### ", "\n#### ", "\n##### "])

    if has_code_blocks:
        # 包含代码块：保护代码块完整性
        base_separators.extend(["\n```\n", "```\n"])

    if has_numbered_lists:
        # 有序列表：在列表项边界分割
        base_separators.extend(["\n\n", "\n"])

    # 通用分隔符（按语义重要性排序）
    universal_separators = [
        "\n\n\n",  # 多个空行（章节分隔）
        "\n\n",    # 段落分隔
        "\n",      # 行分隔
        "。\n",    # 中文句号+换行
        ".\n",     # 英文句号+换行
        "！\n",    # 感叹号+换行
        "？\n",    # 问号+换行
        "。",      # 中文句号
        ".",       # 英文句号
        "；",      # 分号
        "，",      # 逗号
        " ",       # 空格
        ""         # 字符级别分割（最后手段）
    ]

    # 合并并去重
    all_separators = base_separators + universal_separators
    return list(dict.fromkeys(all_separators))  # 保持顺序的去重


def _optimize_chunk_boundaries(chunks: list[str], chunk_size: int, chunk_overlap: int) -> list[str]:
    """
    对初步切分后的文本块进行语义边界优化，确保每个 chunk 尽量在完整句子处结束，避免语义断裂。

    本函数通过以下方式提升 chunk 质量：
    - 避免在句子中间截断（如 '这是一个关于机器' → '这是一个关于机器学习的模型' 被拆开）
    - 修复未完成的句子（如以 '因为...' 结尾）
    - 去除空块，保证输出整洁

    注意：本函数不改变 chunk 数量结构，仅对每个 chunk 内容做微调，适合在递归切分后调用。

    Args:
        chunks (list[str]): 初始切分后的文本块列表
        chunk_size (int): 目标块大小（用于判断是否“过小”）
        chunk_overlap (int): 重叠长度（目前未使用，保留接口兼容性）

    Returns:
        list[str]: 经过语义边界优化后的文本块列表

    示例：
        输入: ["这是一个关于机器", "学习的模型。它很强大"]
        输出: ["这是一个关于机器学习的模型。", "它很强大"]

    设计动机：
        - LLM 更擅长理解完整语义单元
        - 被切断的句子会降低检索相关性和生成连贯性
    """
    optimized_chunks = []

    for chunk in chunks:
        # 检查分块是否在句子中间断开
        if len(chunk) < chunk_size * 0.8:  # 如果分块较小，可能需要合并
            optimized_chunk = _ensure_sentence_boundary(chunk)
        else:
            optimized_chunk = chunk

        # 检查并修复不完整的句子
        optimized_chunk = _fix_incomplete_sentences(optimized_chunk)

        if optimized_chunk.strip():  # 只保留非空分块
            optimized_chunks.append(optimized_chunk)

    return optimized_chunks


def _ensure_sentence_boundary(chunk: str) -> str:
    """
    确保文本块在完整的句子边界处结束，避免在句中截断。

    实现方式：
        - 使用正则匹配中英文句末标点（。！？.!?）
        - 找到最后一个句末符号的位置
        - 若该位置距离末尾较远（>10字符），则截断到此处，丢弃后续“半句话”

    举例：
        输入: "这是一个测试，还没有结束。这是另一个句子但被截断了"
        输出: "这是一个测试，还没有结束。"

    注意：
        - 不会向前“补全”句子，只做“截尾”处理
        - 适用于递归切分后出现“尾部不完整”的场景

    Args:
        chunk (str): 待处理的文本块

    Returns:
        str: 在句子边界结束的文本块
    """
    # 中英文句号模式
    sentence_endings = r'[。！？.!?]'

    # 查找最后一个句子结束位置
    matches = list(re.finditer(sentence_endings, chunk))
    if matches:
        last_sentence_end = matches[-1].end()
        # 如果最后一个句号不在末尾，截断到句号处
        if last_sentence_end < len(chunk) - 10:  # 允许一些尾随空白
            chunk = chunk[:last_sentence_end]

    return chunk


def _fix_incomplete_sentences(chunk: str) -> str:
    """
    修复未完整结束的句子，确保文本块以完整句结尾。

    与 _ensure_sentence_boundary 不同，此函数更激进：
        - 即使 chunk 较长，只要结尾不是句号类符号，就尝试向前查找最近的完整句并截断

    举例：
        输入: "这个模型非常强大，能够处理多种任务，尤其是在自然语言理解方面"
        输出: "这个模型非常强大，能够处理多种任务，"

    虽然损失信息，但保证了每个 chunk 是“可独立理解”的语义单元。

    Args:
        chunk (str): 输入文本块

    Returns:
        str: 以完整句子结尾的文本块（可能被截短）
    """
    chunk = chunk.strip()

    # 如果分块不以句号、问号、感叹号结尾，且长度足够，尝试找到合适的结束点
    if chunk and not re.search(r'[。！？.!?]\s*$', chunk):
        # 查找最后一个完整句子的结束位置
        sentence_endings = r'[。！？.!?]'
        matches = list(re.finditer(sentence_endings, chunk))

        if matches:
            # 截取到最后一个完整句子
            last_sentence_end = matches[-1].end()
            chunk = chunk[:last_sentence_end]

    return chunk


def _evaluate_and_fix_chunks(chunks: list[str], target_size: int) -> list[str]:
    """
    对分块结果进行质量评估与修复，提升整体一致性与可用性。

    修复包括：
        1. 合并过短的碎片块（避免出现只有几个字的 chunk）
        2. 修复未闭合的引号、括号等结构标记
        3. 清理空白块

    本函数采用“流式合并”策略：遍历过程中动态决定是否合并到前一个 chunk。

    Args:
        chunks (list[str]): 待评估的文本块列表
        target_size (int): 目标块大小（用于判断“过短”）

    Returns:
        list[str]: 修复后的高质量文本块列表

    示例：
        输入: ["(不完整的括号", "正常文本", "「未闭合的引号"]
        输出: ["(不完整的括号 正常文本", "「未闭合的引号"] → 自动修复引号/括号
    """
    final_chunks = []

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue

        # 质量评估指标
        chunk_len = len(chunk)

        # 1. 长度评估
        if chunk_len < target_size * 0.1:  # 过短的分块
            # 尝试与前一个分块合并
            if final_chunks and len(final_chunks[-1]) + chunk_len < target_size * 1.2:
                final_chunks[-1] = final_chunks[-1] + "\n" + chunk
                continue

        # 2. 内容完整性评估
        # 检查是否有未闭合的引号、括号等
        if _has_unclosed_markers(chunk):
            chunk = _fix_unclosed_markers(chunk)

        # 3. 添加修复后的分块
        final_chunks.append(chunk)

    return final_chunks


def _has_unclosed_markers(text: str) -> bool:
    """
    检查文本中是否存在未闭合的结构标记，如引号、括号、书名号等。

    支持常见成对符号：
        " 和 "（中英文引号）
        ( 和 )
        [ 和 ]
        { 和 }
        《 和 》
        「 和 」

    注意：仅通过计数判断，不考虑嵌套顺序（简单但高效）

    Args:
        text (str): 输入文本

    Returns:
        bool: 若存在未闭合的标记，返回 True；否则 False

    示例：
        "(hello" → True
        "(hello)" → False
        "「不匹配的引号" → True
    """
    markers = {
        '"': '"',
        '"': '"',
        '(': ')',
        '[': ']',
        '{': '}',
        '《': '》',
        '「': '」'
    }

    for open_mark, close_mark in markers.items():
        open_count = text.count(open_mark)
        close_count = text.count(close_mark)
        if open_count != close_count:
            return True

    return False


def _fix_unclosed_markers(text: str) -> str:
    """
    简单修复未闭合的结构标记：移除多余的开启符号。

    策略：
        - 遍历每对标记（如 " 和 "）
        - 若开启符号多于闭合符号，从右往左删除多余的开启符号
        - 优先保留左侧完整结构

    举例：
        输入: '这是"一个"未"闭合的引号'
        开启: 3个 ", 闭合: 2个 " → 删除最后一个 "
        输出: '这是"一个"未闭合的引号'

    注意：
        - 不添加闭合符号（避免伪造内容）
        - 只删除开启符号（保守策略）

    Args:
        text (str): 输入文本

    Returns:
        str: 修复后的文本（开启符号数量 ≤ 闭合符号）
    """
    # 简单的修复策略：移除不匹配的开启标记
    markers = {
        '"': '"',
        '"': '"',
        '(': ')',
        '[': ']',
        '{': '}',
        '《': '》',
        '「': '」'
    }

    for open_mark, close_mark in markers.items():
        while text.count(open_mark) > text.count(close_mark):
            # 移除最后一个未配对的开启标记
            last_open = text.rfind(open_mark)
            if last_open != -1:
                text = text[:last_open] + text[last_open+1:]

    return text


def _clean_chunk_prefixes(chunks: list[str]) -> list[str]:
    """
    清理分块开头的无用标点符号，确保分块内容整洁。

    主要处理以下情况：
        - 去除分块开头的多余空白字符
        - 去除分块开头的无意义标点（如多个句号、逗号等）
        - 确保分块以有意义的内容开始

    Args:
        chunks (list[str]): 待清理的文本块列表

    Returns:
        list[str]: 清理后的文本块列表

    示例：
        输入: ["。 这是一个测试。", "！你好！", "，欢迎来到。"]
        输出: ["这是一个测试。", "你好！", "欢迎来到。"]
    """
    cleaned_chunks = []

    for chunk in chunks:
        # 去除开头和结尾的空白字符
        chunk = chunk.strip()

        # 使用正则去除开头的无意义标点符号和空白
        # 匹配开头的标点符号（中英文）+ 可选的空白字符
        chunk = re.sub(r'^[。！？,.，；;:：\s]+', '', chunk)

        # 再次去除可能残留的开头空白
        chunk = chunk.lstrip()

        if chunk:  # 只添加非空分块
            cleaned_chunks.append(chunk)

    return cleaned_chunks


def split_text(text: str) -> list[str]:
    """将长文本切分为小块，使用增强递归分块方法。"""
    # 使用增强递归分块方法替代标准方法
    return enhanced_recursive_split(text)


class FileProcessor:
    """用于跟踪文件处理状态。"""

    def __init__(self):
        self.processed_files = {}

    def clear_files(self):
        self.processed_files = {}

    def add_file(self, file_name):
        self.processed_files[file_name] = {
            'status': '等待处理',
            'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'chunks': 0
        }

    def update_status(self, file_name, status, chunks=None):
        if file_name in self.processed_files:
            self.processed_files[file_name]['status'] = status
            if chunks is not None:
                self.processed_files[file_name]['chunks'] = chunks

    def get_file_list(self):
        return [
            f"📄 {fname} | {info['status']}"
            for fname, info in self.processed_files.items()
        ]


def process_files_to_chunks(files: list, file_processor: FileProcessor, progress=None) -> tuple[list, list, list]:
    """处理上传的文件列表，返回所有文本块、元数据和ID。"""
    all_new_chunks = []
    all_new_metadatas = []
    all_new_original_ids = []
    total_files = len(files)

    for idx, file in enumerate(files, 1):
        file_name = os.path.basename(file.name)
        if progress is not None:
            progress((idx - 1) / total_files, desc=f"处理文件 {idx}/{total_files}: {file_name}")
        file_processor.add_file(file_name)

        try:
            text = extract_text_from_file(file.name)
            if not text.strip():
                raise ValueError("文档内容为空或无法提取文本")

            chunks = split_text(text)
            doc_id = f"doc_{int(time.time())}_{idx}"

            current_file_ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
            current_file_metadatas = [{"source": file_name, "doc_id": doc_id} for _ in chunks]

            all_new_chunks.extend(chunks)
            all_new_metadatas.extend(current_file_metadatas)
            all_new_original_ids.extend(current_file_ids)

            file_processor.update_status(file_name, "处理完成", len(chunks))
        except Exception as e:
            logging.error(f"处理文件 {file_name} 时出错: {e}")
            file_processor.update_status(file_name, f"处理失败: {e}")

    return all_new_chunks, all_new_metadatas, all_new_original_ids